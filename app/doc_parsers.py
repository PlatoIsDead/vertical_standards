"""
app/doc_parsers.py — парсеры документов для ингеста (доработка №4).

Диспетчер parse_file: docx (scripts/parse_standards), md (перенос из bitrix_bot),
txt (чанки по ~450 слов), pptx (слайд = чанк), pdf (страница = чанк).
Парсеры возвращают «сырые» чанки {text, heading, section} — роли/audience/doc_name
навешивает process_new_document. Импорты форматных библиотек ленивые: бот не должен
падать на старте, если какая-то из них не установлена.
"""
import os
import re
import sys

SUPPORTED_EXTS = ("docx", "md", "pptx", "txt", "pdf")

_SCRIPTS_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "scripts"))


def parse_file(path: str, ext: str, file_name: str) -> list[dict]:
    """Диспетчер по расширению. Неизвестный формат → [] (файл пропускается)."""
    if ext == "docx":
        if _SCRIPTS_DIR not in sys.path:
            sys.path.insert(0, _SCRIPTS_DIR)
        from parse_standards import parse_docx  # noqa: PLC0415
        chunks = parse_docx(path)
        # Vision-ингест (21.08): SmartArt-схемы и скриншоты → текстовые чанки
        from app.media_ingest import docx_media_chunks  # noqa: PLC0415
        return chunks + docx_media_chunks(path, file_name)
    if ext == "md":
        return parse_md(path)
    if ext == "txt":
        return parse_txt(path, file_name)
    if ext == "pptx":
        return parse_pptx(path, file_name)
    if ext == "pdf":
        return parse_pdf(path, file_name)
    print(f"[doc_parsers] Unsupported ext: .{ext} ({file_name})")
    return []


def parse_md(filepath: str) -> list[dict]:
    """Minimal .md parser: split by headings. (Перенос из bitrix_bot без изменений.)"""
    with open(filepath, "r", encoding="utf-8") as f:
        content = f.read()

    chunks_out = []
    for section in re.split(r"\n(?=#{1,3} )", content):
        lines = section.strip().splitlines()
        if not lines:
            continue
        m = re.match(r"^#{1,3}\s+(.+)", lines[0])
        heading = m.group(1).strip() if m else "Общее"
        body = " ".join(lines[1:]).strip()
        if len(body) > 50:
            chunks_out.append({"text": body, "heading": heading, "section": "Общее"})
    return chunks_out


def parse_txt(path: str, file_name: str) -> list[dict]:
    """Чанки по ~450 слов (конвенция MVP 400–500), heading/section = имя файла."""
    with open(path, "r", encoding="utf-8") as f:
        words = f.read().split()
    out = []
    for i in range(0, len(words), 450):
        body = " ".join(words[i:i + 450])
        if len(body) > 50:
            out.append({"text": body, "heading": file_name, "section": file_name})
    return out


def parse_pptx(path: str, file_name: str) -> list[dict]:
    """Слайд с текстом = чанк; heading = титул слайда или «Слайд N»."""
    from pptx import Presentation  # noqa: PLC0415

    prs = Presentation(path)
    out = []
    for n, slide in enumerate(prs.slides, 1):
        texts = [sh.text_frame.text.strip() for sh in slide.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()]
        body = re.sub(r"\s+", " ", " ".join(texts)).strip()
        if len(body) < 30:  # пустые/мусорные слайды
            continue
        title = ""
        try:
            if slide.shapes.title is not None:
                title = (slide.shapes.title.text or "").strip()
        except (AttributeError, KeyError):
            title = ""
        out.append({"text": body, "heading": title or f"Слайд {n}",
                    "section": file_name})
    return out


def parse_pdf(path: str, file_name: str) -> list[dict]:
    """Страница с текстом = чанк. extract_text() у сканов возвращает "" — пропуск."""
    from pypdf import PdfReader  # noqa: PLC0415

    reader = PdfReader(path)
    out = []
    for n, page in enumerate(reader.pages, 1):
        body = re.sub(r"\s+", " ", (page.extract_text() or "")).strip()
        if len(body) < 30:
            continue
        out.append({"text": body, "heading": f"Стр. {n}", "section": file_name})
    return out
