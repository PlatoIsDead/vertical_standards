"""doc_parsers: txt/pptx/pdf/md — без сети."""
import app.doc_parsers as dp


def test_txt_chunking(tmp_path):
    p = tmp_path / "стандарт.txt"
    p.write_text(" ".join(f"слово{i}" for i in range(1000)), encoding="utf-8")
    chunks = dp.parse_txt(str(p), "стандарт.txt")
    sizes = [len(c["text"].split()) for c in chunks]
    assert sizes == [450, 450, 100]
    assert all(c["heading"] == "стандарт.txt" for c in chunks)
    assert all(c["section"] == "стандарт.txt" for c in chunks)


def test_txt_short_tail_dropped(tmp_path):
    p = tmp_path / "к.txt"
    p.write_text("мало слов", encoding="utf-8")  # < 50 символов
    assert dp.parse_txt(str(p), "к.txt") == []


def test_pptx_slides(tmp_path):
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # титул + контент
    slide.shapes.title.text = "Заголовок слайда"
    slide.placeholders[1].text = "Достаточно длинный текст тела слайда для чанка."
    prs.slides.add_slide(prs.slide_layouts[6])  # пустой
    path = tmp_path / "презентация.pptx"
    prs.save(str(path))

    chunks = dp.parse_pptx(str(path), "презентация.pptx")
    assert len(chunks) == 1
    assert chunks[0]["heading"] == "Заголовок слайда"
    assert "текст тела" in chunks[0]["text"]
    assert chunks[0]["section"] == "презентация.pptx"


def test_pdf_blank_no_crash(tmp_path):
    from pypdf import PdfWriter
    w = PdfWriter()
    w.add_blank_page(width=200, height=200)
    path = tmp_path / "пустой.pdf"
    with open(path, "wb") as f:
        w.write(f)
    assert dp.parse_pdf(str(path), "пустой.pdf") == []


def test_md_moved_parser(tmp_path):
    p = tmp_path / "док.md"
    p.write_text("# Один\n" + "текст " * 20 + "\n## Два\nкоротко", encoding="utf-8")
    chunks = dp.parse_md(str(p))
    assert len(chunks) == 1
    assert chunks[0]["heading"] == "Один"


def test_parse_file_dispatch_and_unknown(tmp_path):
    p = tmp_path / "а.txt"
    p.write_text("слово " * 100, encoding="utf-8")
    assert dp.parse_file(str(p), "txt", "а.txt")
    assert dp.parse_file(str(p), "exe", "а.exe") == []
    assert "txt" in dp.SUPPORTED_EXTS and "pptx" in dp.SUPPORTED_EXTS
